from pptx import Presentation, slide, table
import pandas as pd
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.base import BaseShape
import re
import math


def rgb_to_colour_name(rgb: tuple[int, int, int]) -> str:
    """
    Finds the closest match (green, yellow, or red) by choosing the color with
    the smallest Euclidean distance in rgb value.

    === parameters ===
    rgb: A tuple containing specified rgb values

    === return ===
    example:
    >>> rgb_to_colour_name((253,217,102))
    'Yellow'
    >>> rgb_to_colour_name((41,175,140))
    'Green'
    """
    green = (147, 196, 125)
    yellow = (255, 229, 153)
    red = (213, 97, 97)
    dists = []
    for color in [green, yellow, red]:
        dists.append(math.dist(color, rgb))
    closest = min(dists)
    if closest == dists[0]:
        return "Green"
    elif closest == dists[1]:
        return "Yellow"
    return "Red"


def get_top_dist(shape: BaseShape) -> float:
    """
    Returns distance from top edge of shape to top edge of slideshow in cm
    """
    return shape.top / 360000


def get_left_dist(shape: BaseShape) -> float:
    """
    Returns distance from left edge of shape to left edge of slideshow in cm
    """
    return shape.left / 360000


def is_overlapping(shape: BaseShape, shapes_lst: list) -> int:
    """
    Return index of shape in shapes_lst that under-laps with shape. Return -1
    if there is no such shape.
    Note: shape and all shapes in slide_shapes are Auto Shapes (OKR progress
    indicator shapes).
    """
    for i in range(0, len(shapes_lst)):
        shape_dist = get_top_dist(shape)
        shape_i_dist = get_top_dist(shapes_lst[i])
        if shape_i_dist - 0.3 <= shape_dist  <= shape_i_dist + 0.3:
            return i
    return -1


def get_auto_shapes(slide: slide.Slide) -> list:
    """
    Return all Auto Shapes (OKR progress indicator shapes) on a slide. For any
    overlapping shapes, only keep the topmost shape.
    """
    auto_shapes = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            i = is_overlapping(shape, auto_shapes)
            if i >= 0:
                auto_shapes.pop(i)
            auto_shapes.append(shape)
    return auto_shapes


def sort_auto_shapes(shapes_lst: list) -> list:
    """
    Sort shapes in the vertical order that they are arranged in on the slide.
    Also remove any auto shapes that are not part of the general column created
    by the progress indicator shapes.
    """
    d = {}
    for shape in shapes_lst:
        d[get_top_dist(shape)] = shape
    sorted = []
    for i in range(0, len(d)): # Sort the shapes in the order of the least distance from top edge of slide
        smallest = min(d.keys())
        sorted.append(d[smallest])
        d.pop(smallest)

    sorted2 = [] # Next, only keep shapes that are within the general column
    left_dist = []
    for shape in sorted:
        left_dist.append(get_left_dist(shape))
    avg_left_dist = sum(left_dist) / len(left_dist)
    for i in range(0, len(left_dist)):
        if avg_left_dist - 2 <= left_dist[i] <= avg_left_dist + 2:
            sorted2.append(sorted[i])

    return sorted2


def get_shape_fill_colours(shapes_lst: list) -> list[str]:
    """
    Return a list of corresponding fill colours.
    """
    colors = []
    for shape in shapes_lst:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            color_obj = shape.fill.fore_color
            if color_obj.type == 1:
                colors.append(rgb_to_colour_name(color_obj.rgb))
            elif color_obj.type == 2:
                if color_obj.theme_color == 9:
                    colors.append("Green")
                elif color_obj.theme_color == 10:
                    colors.append("Yellow")
                else:
                    colors.append("Red")
            else:
                colors.append(None) # indicates that fill color was not rgb, nor a theme color
    return colors


def get_progress_colours_from_slide(slide: slide.Slide) -> list[str]:
    """
    Returns progress indicator colours from an OKR table slide.
    Note: When shapes are overlapping, only the top most shapes are included.
    Shapes are also sorted in the order they appear from top to bottom. Any
    shapes outside the general column are removed as outliers. Progress colours
    are lastly retrieved through shape fill.
    """
    shapes = get_auto_shapes(slide)
    sorted_shapes = sort_auto_shapes(shapes)
    return get_shape_fill_colours(sorted_shapes)


def get_meeting_date(p: Presentation) -> str:
    """
    Returns the meeting date as a string from the title slide of the
    presentation.

    === parameters ===
    p: Presentation Object

    === return ===
    example:
    "May 22, 2025"
    """
    date_pattern = r"[A-Za-z]+\s*\d{1,2}\s*[,]*\s*\d{4}"
    for shape in p.slides[0].shapes:
        if shape.has_text_frame:
            text = shape.text.strip()
            match = re.search(date_pattern, text)
            if match:
                return text
    return "No Date Found"


def get_department_names(p: Presentation, ignore_hidden: bool=False) -> list[tuple[str, int]]:
    """
    This function returns a list of tuples containing all the different department
    names and the slide number of the title slide where they were found at.
    The title slide is ideal to extract the team name because there are not many
    shapes to search through, and we can use the slide number as reference to
    determine which row of data belongs to which department in get_data_from_tables().
    Hint: Try using 'Presented by ...' to determine a title slide in regex.

    === parameters ===
    p: Presentation Object

    === return ===
    example:
    [(Product Experience Team, 4), (Marcom Team, 8), ..., (Finance Team, 52)]
    """
    dept_slide_pattern = r"presented\s*by\s*"
    team_name_pattern = r".+team"
    lst = []
    slide_num = 0
    for slide in p.slides:
        slide_num += 1
        if ignore_hidden and is_slide_hidden(slide):
            continue
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text.strip()
                match = re.search(dept_slide_pattern, text.lower()) # First, we match the title slide
                if match:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            text = shape.text.strip()
                            match2 = re.search(team_name_pattern, text.lower()) # Next, we match the team name string
                            if match2:
                                lst.append((text, slide_num))
    return lst


def add_dept_names(okr_table_data: list[list] , dept_names: list[tuple[str, int]]) -> list[list] :
    """
    Add the respective department names to the OKR table data list extracted by
    get_data_from_tables() and return the new okr table data list.

    === parameters ===
    okr_table_data: nested list of data extract from all okr tables
    dept_names: list of tuples containing unique department names and their respective title slide numbers
    """
    for i in range(0, len(okr_table_data)):
        for j in range(0, len(dept_names)):
            if j == len(dept_names) - 1 and dept_names[j][1] < okr_table_data[i][0]: # The last team
                okr_table_data[i].append(dept_names[j][0])
            elif dept_names[j][1] < okr_table_data[i][0] < dept_names[j+1][1]:
                okr_table_data[i].append(dept_names[j][0])
                continue
    return okr_table_data


def get_member_count(p: Presentation, ignore_hidden: bool=True) -> list[tuple]:
    """
    Return a list of tuples containing:
    (slide number, total number of employees) for each department
    in the order that the departments appear in the presentation.

    === parameters ===
    p: Presentation Object

    === return ===
    [(5, 38), (9, 25), ..., (53, 28)]
    """
    members_pattern = r'total\s*members*:?\s*\D*\d+'
    slide_num = 0
    lst = []
    for slide in p.slides:
        slide_num += 1
        if ignore_hidden and is_slide_hidden(slide):
            continue
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text.strip().lower()
                match = re.match(members_pattern, text)
                if match:
                    total_mem = (match.group().split(' ')[-1])
                    total_mem = int(re.sub('[^0-9]','', total_mem))
                    lst.append((slide_num, total_mem))
    return lst


def add_member_count(okr_table_data: list[list] , member_counts: list[tuple]) -> list[list[str]] :
    """
    Add the respective member count to the OKR table data list extracted by
    get_data_from_tables() and return the new okr table data list.

    === parameters ===
    okr_table_data: nested list of data extract from all okr tables
    member_count: list of tuples containing unique department names and their respective title slide numbers
    """
    for i in range(0, len(okr_table_data)):
        for j in range(0, len(member_counts)):
            if j == len(member_counts) - 1 and member_counts[j][0] < okr_table_data[i][0]: # The last team
                okr_table_data[i].append(member_counts[j][1])
            elif member_counts[j][0] < okr_table_data[i][0] < member_counts[j+1][0]:
                okr_table_data[i].append(member_counts[j][1])
                continue
    return okr_table_data


def is_okr_table(t: table.Table) -> bool:
    """
    Returns whether a table is an OKR table or not

    === parameters ===
    t: Table Object
    """
    header = t.rows[0].cells
    header_match = r'okrs?.*projects?.*owners?.*stakeholders?.*status.*deadlines?.*'
    if len(header) == 6:
        header_str = ''
        for cell in header:
            header_str += cell.text
        match = re.search(header_match, header_str, re.IGNORECASE)
        if match:
            return True
        return False
    return False


def is_slide_hidden(s: slide.Slide) -> bool:
    """
    Returns whether a slide is hidden or not.

    === parameters ===
    s: Slide Object
    """
    return s._element.get("show") == "0"


def get_data_from_tables(p: Presentation, ignore_hidden: bool=False) -> list[list[str]]:
    """
    Returns a nested list containing entries from all the tables in a
    presentation, ignoring specified slide numbers.

    === parameters ==
    p: Presentation Object
    ignore_hidden: True or False, based on whether the user would like to ignore
    OKR tables on hidden slides or not

    === return ===
    [[SlideNumber, OKRs, Project, Owner, Stakeholders, Status, Deadline, ProgressColor],
    [[...,...,...,...,...,...,...,...],
    [[...,...,...,...,...,...,...,...],
    [[...,...,...,...,...,...,...,...]]
    """
    lst = []
    slide_num = 0
    for slide in p.slides: # loop through all the slides
        slide_num += 1
        if ignore_hidden and is_slide_hidden(slide): # ignore hidden slides
            continue
        for shape in slide.shapes:
            if shape.has_table and is_okr_table(shape.table): # detect okr tables
                prev_row = []
                progress_colors = get_progress_colours_from_slide(slide)
                row_num = 0
                for row in shape.table.rows:
                    row_num += 1
                    if row_num == 1: # ignore the header row
                        continue
                    r = [slide_num]
                    for i in range(0, len(row.cells)): # loop through the number of data points we will extract
                        cell = row.cells[i]
                        if cell.text == '': # if the text in a cell is empty, then fetch it from the same cell in the previous row
                            r.append(prev_row[i + 1])
                        else:
                            r.append(cell.text)
                    prev_row = r # update previous row to current row
                    r.append(progress_colors[row_num - 2]) # insert corresponding color value at the end
                    lst.append(r) # add the row to our output list
    return lst


def print_table_summary(p: Presentation()):
    """
    This function is helpful in testing our table search algorithm and
    catching tables that were incorrectly included.
    Print a summary containing the number of total tables and number of OKR tables
    in the presentation. For every table found, indicate whether it is an OKR table
    or not.

    === parameters ===
    p: Presentation Object
    """
    num_match = 0
    num_tables = 0
    slide_num = 0
    for slide in p.slides:
        slide_num += 1
        for shape in slide.shapes:
            if shape.has_table:
                num_tables += 1
                if is_okr_table(shape.table):
                    num_match += 1
                    print(f"OKR Table on slide {slide_num}")
                else:
                    print(f"Not an OKR Table on slide {slide_num}")
    print(f"\n{num_tables} total tables found, with {num_match} matches for an OKR table.")


if __name__ == '__main__':
    pres = Presentation('OKR_Testing.pptx')
    ign_hidden = True
    dept_names = get_department_names(pres, ign_hidden)
    data = get_data_from_tables(pres, ign_hidden)
    data = add_dept_names(data, dept_names)
    data = add_member_count(data, get_member_count(pres, ign_hidden))

    df = pd.DataFrame(data, columns=['SlideNumber', 'OKRs', 'Projects', 'Owner', 'Stakeholders', 'Status', 'Deadline', 'ProgressColor', 'Team', 'TotalMembers'])
    df['MeetingDate'] = [get_meeting_date(pres)] * (len(data)) # Adds a column for meeting date
    df.to_excel('okr_table_data.xlsx', index=False)

    print_table_summary(pres)
    print(f"{df.shape[0]} total OKR table entries successfully extracted.")

